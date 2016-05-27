# -----------------------------------------------------
# Name: sysmap_pptx.py
# Function: Read input file to calc & write to PPTX file.
# Usage: $ python sysmap_pptx.py <input txt file> <output pptx file>[CR]
# Output: pptx format file
# -----------------------------------------------------
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Mm, Pt
import sys
import string

HEADER = "Any["
FOOTER = "]"
SPLITTER = '--- SPLITTER ---'   # Must same with sysmap_jl.jl def.
SLIDE_TYPE = 6  # 6: Blank slide. 5: One title slide
TITLE = "Technology System Map"  # Only used in slide_title.

print("=== sysmap_pptx: Start.")

if __name__ == "__main__":

    prs = Presentation()
    #--- Default slide width
    prs.slide_width = 9144000
    # prs.slide_height = 6858000  # 4:3
    prs.slide_height = 5143500  # 16:9

    #--- create one blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_TYPE])
    # if use the following title, set slide 5 instead of 6(blank)
    # in the upper code line.
    # title = slide.shapes.title
    # title.text = 'Technology System Map'
    # title.text = TITLE

    
    #--- Set each item to the slide
    fn = sys.argv[1]   # input position data file.
    fd = open(fn, 'r')
    fo = sys.argv[2]
    # print "fn :  %s" % fn
    # print "fo:  %s" % fo

    while 1:
        s = fd.readline()
        if not s:
            left = top = width = height = 1000
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = " End of System Map."
            break
        s = s[:-1]       #--- Remove \n char for compare
        #---------------------------------
        #--- Add slide if new Appl group.
        #---------------------------------
        if s == SPLITTER:
            slide = prs.slides.add_slide(prs.slide_layouts[SLIDE_TYPE])
            # Used in case of slide type is Title
            # title = slide.shapes.title
            # title.text = TITLE
            continue
        #---------------------------------
        #--- Remove non-needed char.
        #---------------------------------
        s = string.replace(s, HEADER, "") # del head of str
        s = string.replace(s, FOOTER, "") # del tail of str
        parts = s.split(",")
        label = parts[0]
        label = string.replace(label, "\"", "") # del double quotation
        #---------------------------------
        #--- Calc obj position in slide.
        #---------------------------------
        x = Mm(int(float(parts[1])))
        y = Mm(int(float(parts[2])))
        w = Mm(int(float(parts[3])))
        h = Mm(int(float(parts[4])))
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        shape.text = '%s' % label
        # print "Shape.text %s" % shape.text
        # print(shape.id, shape.text.encode('utf-8'), shape.name, shape.left, shape.top, shape.width, shape.height)

    fd.close()
    prs.save(fo)

    print("=== sysmap_pptx: End. Output file is %s" % fo)

# ----------------------------------------------
# End of File  
# ----------------------------------------------
